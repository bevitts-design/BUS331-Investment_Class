#!/usr/bin/env python3
"""Create a path-safe, slide-level audit of the BUS331 Chapters 1-4 source PPTX."""

from __future__ import annotations

import argparse
import json
from collections import Counter
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def clean(value: str | None) -> str:
    lines = []
    for line in (value or "").replace("\v", "\n").splitlines():
        normalized = " ".join(line.split())
        if normalized:
            lines.append(normalized)
    return "\n".join(lines)


def shape_metadata(shape: Any) -> dict[str, str]:
    metadata = {"name": clean(getattr(shape, "name", ""))}
    try:
        c_nv_pr = shape._element.xpath(".//p:cNvPr")[0]
        metadata["title"] = clean(c_nv_pr.get("title"))
        metadata["description"] = clean(c_nv_pr.get("descr"))
    except (IndexError, AttributeError):
        metadata["title"] = ""
        metadata["description"] = ""
    return metadata


def walk_shapes(shapes: Any) -> list[Any]:
    found = []
    for shape in shapes:
        found.append(shape)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            found.extend(walk_shapes(shape.shapes))
    return found


def slide_title(slide: Any, text_blocks: list[str], slide_number: int) -> str:
    if slide.shapes.title is not None:
        title = clean(slide.shapes.title.text).replace("\n", " ")
        if title and title != str(slide_number):
            return title
    for value in text_blocks:
        if value and value != str(slide_number) and not value.startswith("©"):
            return value.replace("\n", " ")[:180]
    return f"Visual source slide {slide_number}"


def audit(source: Path) -> dict[str, Any]:
    presentation = Presentation(source)
    slides = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        shapes = walk_shapes(slide.shapes)
        text_blocks: list[str] = []
        tables: list[dict[str, Any]] = []
        visuals: list[dict[str, str]] = []
        type_counts: Counter[str] = Counter()
        chart_count = 0

        for shape in shapes:
            type_name = getattr(shape.shape_type, "name", str(shape.shape_type))
            type_counts[type_name] += 1
            if getattr(shape, "has_text_frame", False):
                value = clean(shape.text)
                if value:
                    text_blocks.append(value)
            if getattr(shape, "has_table", False):
                rows = [[clean(cell.text) for cell in row.cells] for row in shape.table.rows]
                tables.append({"rows": rows, "rowCount": len(rows), "columnCount": len(rows[0]) if rows else 0})
            if getattr(shape, "has_chart", False):
                chart_count += 1
            has_image_fill = shape.shape_type != MSO_SHAPE_TYPE.GROUP and bool(
                shape._element.xpath(".//a:blip")
            )
            if has_image_fill or shape.shape_type in {
                MSO_SHAPE_TYPE.PICTURE,
                MSO_SHAPE_TYPE.LINKED_PICTURE,
                MSO_SHAPE_TYPE.MEDIA,
                MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
                MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
            }:
                visuals.append(shape_metadata(shape))

        note = ""
        try:
            note = clean(slide.notes_slide.notes_text_frame.text)
        except (AttributeError, ValueError):
            pass

        slides.append(
            {
                "slide": slide_number,
                "title": slide_title(slide, text_blocks, slide_number),
                "textBlocks": text_blocks,
                "tables": tables,
                "visuals": visuals,
                "chartCount": chart_count,
                "shapeCounts": dict(sorted(type_counts.items())),
                "speakerNote": note,
            }
        )

    return {
        "sourceFile": source.name,
        "sourceSlideCount": len(slides),
        "slides": slides,
    }


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("source", type=Path)
    parser.add_argument("--output", required=True, type=Path)
    args = parser.parse_args()
    result = audit(args.source)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    args.output.write_text(json.dumps(result, indent=2, ensure_ascii=False), encoding="utf-8")
    print(f"Audited {result['sourceSlideCount']} slides to {args.output}")


if __name__ == "__main__":
    main()
